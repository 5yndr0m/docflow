"""
Document text extraction.

Each parser returns a list of page segments:
    [{"text": str, "page_number": int}, ...]

All parsers are synchronous — call them via asyncio.to_thread() from async code.
"""
from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Characters per synthetic "page" for formats with no native page concept
_SYNTHETIC_PAGE_SIZE = 3000


def parse_pdf(file_path: str) -> list[dict]:
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages.append({"text": text, "page_number": i + 1})
    doc.close()
    return pages


def parse_docx(file_path: str) -> list[dict]:
    from docx import Document

    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []

    full_text = "\n".join(paragraphs)
    segments = []
    for i in range(0, len(full_text), _SYNTHETIC_PAGE_SIZE):
        chunk = full_text[i : i + _SYNTHETIC_PAGE_SIZE].strip()
        if chunk:
            segments.append({"text": chunk, "page_number": i // _SYNTHETIC_PAGE_SIZE + 1})
    return segments


def parse_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        if texts:
            pages.append({"text": "\n".join(texts), "page_number": i + 1})
    return pages


def parse_xlsx(file_path: str) -> list[dict]:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    segments = []
    seg_idx = 1

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                rows.append(row_text)

        # Group every 50 rows into one segment
        for i in range(0, len(rows), 50):
            text = f"Sheet: {sheet_name}\n" + "\n".join(rows[i : i + 50])
            segments.append({"text": text, "page_number": seg_idx})
            seg_idx += 1

    wb.close()
    return segments


def parse_txt(file_path: str) -> list[dict]:
    with open(file_path, encoding="utf-8", errors="replace") as f:
        content = f.read()

    if not content.strip():
        return []

    segments = []
    for i in range(0, len(content), _SYNTHETIC_PAGE_SIZE):
        text = content[i : i + _SYNTHETIC_PAGE_SIZE].strip()
        if text:
            segments.append({"text": text, "page_number": i // _SYNTHETIC_PAGE_SIZE + 1})
    return segments


_PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "xlsx": parse_xlsx,
    "txt": parse_txt,
}


def parse_document(file_path: str, file_type: str) -> list[dict]:
    """Dispatch to the correct parser based on file_type.

    Args:
        file_path: Absolute path to the file on disk.
        file_type: Lowercase extension without dot (e.g. "pdf", "docx").

    Returns:
        List of {text, page_number} dicts. Empty list if no content extracted.
    """
    parser = _PARSERS.get(file_type)
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")

    try:
        segments = parser(file_path)
    except Exception as exc:
        logger.error("Failed to parse %s (%s): %s", file_path, file_type, exc)
        raise

    # Filter out empty segments
    return [s for s in segments if s.get("text", "").strip()]
